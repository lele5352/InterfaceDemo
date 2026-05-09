"""
生成项目介绍 PPT - InterfaceDemo 零代码接口自动化框架
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

CLR_PRIMARY = RGBColor(0x1A, 0x73, 0xE8)
CLR_SECONDARY = RGBColor(0x34, 0xA8, 0x53)
CLR_ACCENT = RGBColor(0xEA, 0x43, 0x35)
CLR_DARK = RGBColor(0x20, 0x20, 0x20)
CLR_GRAY = RGBColor(0x66, 0x66, 0x66)
CLR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CLR_CODE_BG = RGBColor(0x28, 0x2C, 0x34)
CLR_CODE = RGBColor(0x98, 0xC3, 0x79)
CLR_ORANGE = RGBColor(0xE6, 0x7E, 0x22)


def add_shape_bg(slide, color=CLR_PRIMARY):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(0.08))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def add_page_number(slide, num, total):
    txBox = slide.shapes.add_textbox(Inches(12.3), Inches(7.1), Inches(0.8), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num}/{total}"
    p.font.size = Pt(10)
    p.font.color.rgb = CLR_GRAY
    p.alignment = PP_ALIGN.RIGHT


def set_cell_text(cell, text, bold=False, size=11, color=CLR_DARK, align=PP_ALIGN.LEFT):
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_section_header(slide, title, subtitle=""):
    add_shape_bg(slide, CLR_PRIMARY)
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(1.2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = CLR_PRIMARY
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(16)
        p2.font.color.rgb = CLR_GRAY


def add_code_block(slide, left, top, width, height, code_text, font_size=11):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CLR_CODE_BG
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.12)
    tf.margin_bottom = Inches(0.08)
    p = tf.paragraphs[0]
    p.text = code_text
    p.font.size = Pt(font_size)
    p.font.color.rgb = CLR_CODE
    p.font.name = "Courier New"


def add_info_box(slide, left, top, width, height, text, color=CLR_SECONDARY):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
    shape.line.color.rgb = color
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.color.rgb = CLR_DARK


def add_bullet_slide(slide, title, bullets, note=""):
    add_section_header(slide, title)
    add_shape_bg(slide, CLR_PRIMARY)
    left = Inches(0.8)
    top = Inches(1.8)
    for i, bullet in enumerate(bullets):
        txBox = slide.shapes.add_textbox(left, top + Inches(i * 0.55), Inches(11.5), Inches(0.6))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        if bullet.startswith("##"):
            p.text = bullet[2:].strip()
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = CLR_PRIMARY
        else:
            p.text = bullet
            p.font.size = Pt(14)
            p.font.color.rgb = CLR_DARK
            p.space_after = Pt(4)
    if note:
        add_info_box(slide, Inches(0.8), Inches(6.0), Inches(11.5), Inches(1.2), note)


def add_two_column(slide, title, left_title, left_items, right_title, right_items):
    add_section_header(slide, title)
    add_shape_bg(slide, CLR_PRIMARY)
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(5.8), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = CLR_SECONDARY
    l_top = Inches(2.2)
    for i, item in enumerate(left_items):
        txBox = slide.shapes.add_textbox(Inches(0.7), l_top + Inches(i * 0.5), Inches(5.5), Inches(0.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(11)
        p.font.color.rgb = CLR_DARK
    txBox = slide.shapes.add_textbox(Inches(6.8), Inches(1.6), Inches(5.8), Inches(0.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = CLR_ORANGE
    r_top = Inches(2.2)
    for i, item in enumerate(right_items):
        txBox = slide.shapes.add_textbox(Inches(7.0), r_top + Inches(i * 0.5), Inches(5.5), Inches(0.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(11)
        p.font.color.rgb = CLR_DARK


TOTAL = 30

# ====== Slide 1: 封面 ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
bg_shape.fill.solid()
bg_shape.fill.fore_color.rgb = CLR_PRIMARY
bg_shape.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(1.5), Inches(1.5), Inches(10), Inches(1.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "InterfaceDemo"
p.font.size = Pt(52)
p.font.bold = True
p.font.color.rgb = CLR_WHITE
p.alignment = PP_ALIGN.CENTER

txBox2 = slide.shapes.add_textbox(Inches(1.5), Inches(3.0), Inches(10), Inches(1))
tf2 = txBox2.text_frame
p2 = tf2.paragraphs[0]
p2.text = "零代码接口自动化测试框架"
p2.font.size = Pt(28)
p2.font.color.rgb = CLR_WHITE
p2.alignment = PP_ALIGN.CENTER

txBox3 = slide.shapes.add_textbox(Inches(2), Inches(4.0), Inches(9), Inches(1))
tf3 = txBox3.text_frame
p3 = tf3.paragraphs[0]
p3.text = "Python + Requests + Pytest + Allure + Jenkins"
p3.font.size = Pt(18)
p3.font.color.rgb = RGBColor(0xBB, 0xDE, 0xFB)
p3.alignment = PP_ALIGN.CENTER

txBox4 = slide.shapes.add_textbox(Inches(2), Inches(5.8), Inches(9), Inches(0.5))
tf4 = txBox4.text_frame
p4 = tf4.paragraphs[0]
p4.text = "面向测试同学的零代码解决方案 | 只需维护 YAML 文件即可完成接口自动化"
p4.font.size = Pt(14)
p4.font.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
p4.alignment = PP_ALIGN.CENTER
add_page_number(slide, 1, TOTAL)

# ====== Slide 2: 目录 ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "目录", "CONTENTS")
add_shape_bg(slide, CLR_PRIMARY)
items = [
    ("01", "项目定位与技术栈"), ("02", "项目目录结构一览"),
    ("03", "环境配置入口"), ("04", "YAML 测试用例详解"),
    ("05", "核心模块逐个解析"), ("06", "热加载函数库"),
    ("07", "运行入口与 Pytest 集成"), ("08", "完整执行流程"),
    ("09", "编写你的第一个测试用例"), ("10", "常见问题 FAQ"),
]
for i, (num, title) in enumerate(items):
    row = i // 2
    col = i % 2
    left = Inches(1 + col * 6)
    top = Inches(1.8 + row * 1.1)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = CLR_PRIMARY
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = str(num)
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = CLR_WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    txBox = slide.shapes.add_textbox(left + Inches(0.65), top + Inches(0.05), Inches(5), Inches(0.5))
    tf2 = txBox.text_frame
    tf2.paragraphs[0].text = title
    tf2.paragraphs[0].font.size = Pt(15)
    tf2.paragraphs[0].font.color.rgb = CLR_DARK
add_page_number(slide, 2, TOTAL)

# ====== Slide 3: 项目定位 ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "项目定位", "PROJECT OVERVIEW")
add_shape_bg(slide, CLR_PRIMARY)
add_bullet_slide(slide, "项目定位", [
    "## 这是什么？",
    "「零代码」接口自动化测试框架 —— 测试同学只需要写 YAML 文件，不用写 Python 代码",
    "## 目标用户",
    "具备基础 Python 能力的测试工程师，希望快速搭建接口自动化体系",
    "## 核心理念",
    "测试用例 = YAML 文件（结构化数据，声明式描述）",
    "框架自动解析 YAML → 发送 HTTP 请求 → 提取响应数据 → 断言验证",
    "用「配置」代替「编码」，降低自动化门槛",
    "## 适合场景",
    "RESTful API 回归测试 | 多环境测试 | 数据驱动的大数据量测试",
])

# ====== Slide 4: 技术栈 ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "技术栈", "TECH STACK")
add_shape_bg(slide, CLR_PRIMARY)
rows, cols = 9, 3
table = slide.shapes.add_table(rows, cols, Inches(1.5), Inches(1.8), Inches(10), Inches(5)).table
headers = ["技术", "版本", "用途"]
data = [
    ["Python", "3.9+", "框架开发语言"],
    ["requests", "2.33+", "HTTP 请求客户端（发送 API 请求）"],
    ["pytest", "8.3+", "测试执行引擎（发现并运行用例）"],
    ["allure-pytest", "2.15+", "测试报告生成（美观的 HTML 报告）"],
    ["PyYAML", "6.0+", "解析 YAML 测试用例文件"],
    ["jsonpath", "0.82+", "JSON 响应数据提取（如 $.data.token）"],
    ["Cerberus", "1.3+", "YAML 用例数据格式校验"],
    ["python-dotenv", "1.2+", "加载 run.env 环境变量"],
]
for i, h in enumerate(headers):
    set_cell_text(table.cell(0, i), h, bold=True, size=13, color=CLR_WHITE)
    table.cell(0, i).fill.solid()
    table.cell(0, i).fill.fore_color.rgb = CLR_PRIMARY
for r, rd in enumerate(data):
    for c, ct in enumerate(rd):
        set_cell_text(table.cell(r+1, c), ct, size=12)
        if r % 2 == 0:
            table.cell(r+1, c).fill.solid()
            table.cell(r+1, c).fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xFF)
add_page_number(slide, 4, TOTAL)

# ====== Slide 5: 架构图 ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "整体架构", "ARCHITECTURE")
add_shape_bg(slide, CLR_PRIMARY)

y_start = Inches(1.6)
layer_h = Inches(0.5)
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), y_start, Inches(12.3), layer_h)
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
shape.line.color.rgb = CLR_SECONDARY
shape.line.width = Pt(2)
tf = shape.text_frame
tf.paragraphs[0].text = "用户层（零代码）— 只需维护 YAML 文件和配置文件，无需编写 Python 代码"
tf.paragraphs[0].font.size = Pt(13)
tf.paragraphs[0].font.color.rgb = CLR_DARK

y2 = y_start + layer_h + Inches(0.15)
shape2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), y2, Inches(12.3), Inches(1.8))
shape2.fill.solid()
shape2.fill.fore_color.rgb = RGBColor(0xE3, 0xF2, 0xFD)
shape2.line.color.rgb = CLR_PRIMARY
shape2.line.width = Pt(2)
tf2 = shape2.text_frame
tf2.word_wrap = True
tf2.paragraphs[0].text = "框架核心层（commons/）— 自动处理一切"
tf2.paragraphs[0].font.size = Pt(13)
tf2.paragraphs[0].font.color.rgb = CLR_DARK
tf2.paragraphs[0].font.bold = True
modules = ["yaml_util → 读取 YAML", "model_util → 校验格式", "request_util → 发送请求 + 变量替换",
           "extract_util → 提取响应数据", "assert_util → 断言验证", "ddt_util → 数据驱动"]
for i, mod in enumerate(modules):
    col = i % 3
    row = i // 3
    txBox = slide.shapes.add_textbox(Inches(0.8 + col * 4), y2 + Inches(0.4 + row * 0.55), Inches(3.8), Inches(0.5))
    p = txBox.text_frame.paragraphs[0]
    p.text = f"  -> {mod}"
    p.font.size = Pt(12)
    p.font.color.rgb = CLR_DARK

y3 = y2 + Inches(1.8) + Inches(0.15)
shape3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), y3, Inches(12.3), layer_h)
shape3.fill.solid()
shape3.fill.fore_color.rgb = RGBColor(0xFD, 0xF2, 0xE9)
shape3.line.color.rgb = CLR_ORANGE
shape3.line.width = Pt(2)
tf3 = shape3.text_frame
tf3.paragraphs[0].text = "基础设施 — Pytest（测试执行）+ Allure（报告生成）+ Jenkins（CI/CD 集成）"
tf3.paragraphs[0].font.size = Pt(13)
tf3.paragraphs[0].font.color.rgb = CLR_DARK

y4 = y3 + layer_h + Inches(0.15)
shape4 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), y4, Inches(12.3), layer_h)
shape4.fill.solid()
shape4.fill.fore_color.rgb = RGBColor(0xFC, 0xE4, 0xEC)
shape4.line.color.rgb = CLR_ACCENT
shape4.line.width = Pt(2)
tf4 = shape4.text_frame
tf4.paragraphs[0].text = "输出 — 控制台日志 + Allure HTML 测试报告 + 变量持久化文件 (extract.yaml)"
tf4.paragraphs[0].font.size = Pt(13)
tf4.paragraphs[0].font.color.rgb = CLR_DARK

add_page_number(slide, 5, TOTAL)

# ====== Slide 6: 目录结构 ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "项目目录结构", "PROJECT STRUCTURE")
add_shape_bg(slide, CLR_PRIMARY)
tree_text = """InterfaceDemo/
  run.py                 运行入口
  debug_talk.py          热加载函数库
  extract.yaml           接口变量持久化
  run.env                环境变量
  pytest.ini             pytest 配置
  commons/               框架核心模块
    settings.py          全局配置单例
    yaml_util.py         YAML 读写工具
    model_util.py        数据模型 + 校验
    request_util.py      HTTP 请求引擎
    extract_util.py      响应数据提取
    assert_util.py       断言验证
    ddt_util.py          数据驱动
    case_util.py         用例发现
    main_util.py         主执行器
    logger.py            日志
  configs/config.yaml    多环境配置
  testcases/             测试用例目录
    conftest.py          fixtures
    test_all_case.py     动态注册引擎
    weather/weather_search.yaml  示例用例
  logs/ /reports/ /temps/  /data/"""
add_code_block(slide, Inches(0.5), Inches(1.6), Inches(12.3), Inches(5.5), tree_text, font_size=11)
add_page_number(slide, 6, TOTAL)

# ====== Slide 7: 环境配置 ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "环境配置入口", "CONFIGURATION")
add_shape_bg(slide, CLR_PRIMARY)

txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(5.5), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "run.env — 环境变量"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = CLR_SECONDARY
code1 = "weather_base_url=https://cy.weather.com.cn\nDB_HOST=localhost\nAPI_KEY=abc123xyz456"
add_code_block(slide, Inches(0.5), Inches(2.2), Inches(5.8), Inches(1.8), code1, font_size=11)

txBox2 = slide.shapes.add_textbox(Inches(7), Inches(1.6), Inches(5.5), Inches(0.5))
tf2 = txBox2.text_frame
p2 = tf2.paragraphs[0]
p2.text = "configs/config.yaml — 多环境"
p2.font.size = Pt(18)
p2.font.bold = True
p2.font.color.rgb = CLR_ORANGE
code2 = "base_urls:\n  dev: http://dev.api.example.com\n  test: https://cy.weather.com.cn\n  prod: https://api.example.com\nheaders:\n  test:\n    Content-Type: application/json\ntimeout: 10"
add_code_block(slide, Inches(7), Inches(2.2), Inches(5.8), Inches(2.8), code2, font_size=11)

note = "YAML 中通过 ${env(weather_base_url)} 引用 run.env\n运行时通过 --env test 切换环境，框架自动读取对应 base_url\ncommons/settings.py 负责加载这两个配置"
add_info_box(slide, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.8), note)
add_page_number(slide, 7, TOTAL)

# ====== Slide 8: YAML 格式 ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "YAML 测试用例格式", "TEST CASE FORMAT")
add_shape_bg(slide, CLR_PRIMARY)
code = """-
  feature: 天气网
  story: 查询当前地天气
  title: 验证天气网查询接口返回成功
  request:
    method: get
    url: ${env(weather_base_url)}/api/v1/promo_pos?
    params:
      ids: 3
    # json: {"key": "value"}   # POST JSON
    # data: {"key": "value"}   # POST 表单
  extract:
    city: [text, '"msg":"(.*?)"', 1]     # 正则提取
    id: [json, "$.data.list[0].id", 0]   # JSONPath
  validate:
    equals:
      断言状态码200: [200, status_code]
    contains:
      断言包含成功: [成功, text]"""
add_code_block(slide, Inches(0.5), Inches(1.6), Inches(8.5), Inches(5.5), code, font_size=11)
keys = [("feature/story/title", "Allure 报告三级标题"),
        ("request (必需)", "HTTP 请求描述"),
        ("extract (可选)", "提取字段用于接口关联"),
        ("validate (可选)", "响应断言验证"),
        ("parametrize (可选)", "数据驱动多组参数")]
for i, (k, d) in enumerate(keys):
    top = Inches(1.7 + i * 1.0)
    txBox = slide.shapes.add_textbox(Inches(9.3), top, Inches(3.5), Inches(0.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = k
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = CLR_PRIMARY
    txBox2 = slide.shapes.add_textbox(Inches(9.3), top + Inches(0.3), Inches(3.5), Inches(0.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = d
    p2.font.size = Pt(11)
    p2.font.color.rgb = CLR_GRAY
add_page_number(slide, 8, TOTAL)

# ====== Slide 9: 变量替换 ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "YAML 变量替换语法", "VARIABLE REFERENCE")
add_shape_bg(slide, CLR_PRIMARY)
add_two_column(slide, "YAML 变量替换语法",
    "变量引用",
    ["${env(KEY)} — 引用 run.env 环境变量\n  例: ${env(weather_base_url)} -> https://cy.weather.com.cn",
     "$ddt{field} — 参数化当前行值\n  例: $ddt{username} -> baili",
     "${read_yaml(key)} — extract.yaml 中的值\n  例: ${read_yaml(city)} -> 广州"],
    "热加载函数",
    ["${random_phone()} -> 19960152366",
     "${random_number(1,100)} -> 42",
     "${uuid_str()} -> bb5db327a21d438c9097bf5cde470055",
     "${timestamp()} -> 1778234480",
     "${current_date()} -> 2026-05-08",
     "${format_time(7)} -> 2026-05-15"])
note = "变量替换在请求发送前由 request_util.py 自动递归完成"
add_info_box(slide, Inches(0.5), Inches(6.0), Inches(12.3), Inches(1.0), note)
add_page_number(slide, 9, TOTAL)

# ====== Slide 10: 数据驱动 ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "数据驱动（参数化）", "DATA-DRIVEN")
add_shape_bg(slide, CLR_PRIMARY)
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(6), Inches(0.5))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "方式一：YAML 内联参数"
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = CLR_PRIMARY
code1 = """-
  title: 登录测试
  request:
    method: get
    url: ${env(base_url)}/login/
    data:
      username: $ddt{username}
      password: $ddt{password}
  parametrize:
    - ["username","password"]
    - ["baili","baili123"]
    - ["admin","admin123"]"""
add_code_block(slide, Inches(0.5), Inches(2.2), Inches(6), Inches(2.8), code1, font_size=10)

txBox2 = slide.shapes.add_textbox(Inches(6.8), Inches(1.6), Inches(6), Inches(0.5))
tf2 = txBox2.text_frame
p2 = tf2.paragraphs[0]
p2.text = "方式二：CSV 文件"
p2.font.size = Pt(16)
p2.font.bold = True
p2.font.color.rgb = CLR_ORANGE
code2 = "# YAML:\nparametrize: data/login.csv\n\n# data/login.csv:\nusername,password,expect\nbaili,baili123,success\nadmin,admin123,success"
add_code_block(slide, Inches(6.8), Inches(2.2), Inches(6), Inches(2.2), code2, font_size=10)
note = "框架自动展开为 pytest 参数化测试，每个数据行 = 一个独立用例"
add_info_box(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.5), note)
add_page_number(slide, 10, TOTAL)

# ====== Slide 11-20: 模块详解 ======
module_slides = [
    ("commons/settings.py", "全局配置单例", [
        "## 作用",
        "整个框架的「中央空调」— 所有模块都通过它获取配置信息",
        "## 核心特性",
        "单例模式 ConfigSingleton，整个框架只有一个实例",
        "自动加载 run.env + configs/config.yaml",
        "提供 get_env() / get_config() / get_base_url() / get_headers() / get_timeout()",
        "set_env(env_name) — 切换环境，自动换 base_url",
        "set_variable / get_variable — 内存变量池（接口间传递数据）",
        "## 使用方式",
        "from commons.settings import config",
        "config.get_base_url()  ->  'https://cy.weather.com.cn'",
    ]),
    ("commons/yaml_util.py", "YAML 读写工具", [
        "## 作用",
        "统一管理所有 YAML 文件的读取和写入",
        "## 核心函数",
        "read_yaml_file(path) — 读取任意 YAML 文件，返回 dict/list",
        "write_yaml_file(path, data) — 写入数据到 YAML 文件",
        "append_to_yaml(path, data) — 追加数据到已有 YAML 文件",
        "clear_yaml_file(path) — 清空 YAML 文件",
        "## extract.yaml 专用",
        "read_extract_yaml() — 读取所有已提取的变量",
        "write_extract_yaml(data) — 追加变量到 extract.yaml",
        "get_extract_value(key) — 获取指定 key 的值",
        "clear_extract_yaml() — 每次测试开始前清空",
    ]),
    ("commons/model_util.py", "数据模型 + 格式校验", [
        "## 作用",
        "定义 YAML 测试用例的数据结构，校验每个字段的合法性",
        "## TestCaseObj 类",
        "把 YAML 字典转为 Python 对象，含 feature/story/title/request/extract/validate 属性",
        "## 关键函数",
        "validate_case_data(case_data) — 校验 YAML 字典是否符合规范",
        "build_test_case_obj(case_data) — 构建 TestCaseObj 对象",
        "verify_yaml(case_data, filename) — 一站式校验 + 构建",
        "## 校验规则",
        "feature/story/title 必须为非空字符串",
        "request 必须含 method (get/post/put/delete) 和 url",
        "校验失败时明确提示：'字段 request: method: unallowed value'",
    ]),
    ("commons/ddt_util.py", "数据驱动引擎", [
        "## 作用",
        "同一 YAML 用例用不同数据「跑多遍」",
        "## 核心函数",
        "load_csv(csv_path) — 读取 CSV，返回字典列表",
        "parse_parametrize(data) — 解析 YAML 内联参数列表",
        "build_parametrize_from_csv(path) — 从 CSV 构建参数化数据",
        "guess_parametrize(raw_data) — 智能识别数据来源",
        "## 执行效果",
        "3 行测试数据 = 3 个独立 pytest 用例，报告中可分别查看",
    ]),
    ("commons/request_util.py", "HTTP 请求引擎", [
        "## 作用",
        "框架的「发动机」— 接收用例描述、发送 HTTP 请求、返回响应",
        "## RequestEngine.send(method, url, **kwargs)",
        "发送前自动递归替换所有变量占位符：",
        "  1) $ddt{field} -> 当前参数化行值",
        "  2) ${env(KEY)} -> 环境变量",
        "  3) ${read_yaml(key)} -> extract.yaml 的值",
        "  4) ${func(args)} -> 热加载函数返回值",
        "## 其他功能",
        "自动管理 Cookie（requests.Session）",
        "自动补全 base_url，设置默认超时和请求头",
    ]),
    ("commons/extract_util.py", "响应数据提取", [
        "## 作用",
        "从 HTTP 响应中提取数据，实现「接口关联」",
        "## ExtractEngine.extract(response, config)",
        "JSONPath 提取: [json, '$.data.token', 0]",
        "  从 JSON 响应体中取 $.data.token，索引 0",
        "正则提取: [text, '\"token\":\"(.*?)\"', 1]",
        "  从响应文本中正则匹配，取分组 1",
        "## 自动保存",
        "提取到的值自动写入 extract.yaml + 内存变量池",
    ]),
    ("commons/assert_util.py", "断言验证", [
        "## 作用",
        "验证接口返回是否符合预期，是测试的「裁判员」",
        "## ResponseValidator.validate(info)",
        "equals 精确匹配: expect == actual",
        "  YAML: 断言状态码200: [200, status_code]",
        "contains 包含匹配: expect in actual",
        "  YAML: 断言包含成功: [成功, text]",
        "## 实际值来源",
        "status_code / text / json / $.jsonpath",
        "## 断言失败",
        "抛出 AssertionError，pytest 自动捕获",
        "错误信息含：描述、期望值、实际值",
    ]),
    ("commons/main_util.py", "主执行器", [
        "## 作用",
        "编排测试用例的完整生命周期",
        "## execute_test_case(case_obj) 流程",
        "Step 1: 取出用例信息（title/request/extract/validate）",
        "Step 2: 发送请求 — request_engine.send()",
        "Step 3: 提取数据 — extract_engine.extract_response_info()",
        "  结果自动写入 extract.yaml",
        "Step 4: 断言验证 — ResponseValidator.validate()",
        "  失败 -> AssertionError -> pytest 标记失败",
        "Step 5: 记录执行结果到日志",
    ]),
    ("commons/case_util.py", "用例发现与加载", [
        "## 作用",
        "扫描 testcases/ 目录，发现 YAML 文件并解析为 TestCase 对象",
        "## 核心函数",
        "discover_yaml_files() — 递归扫描返回所有 .yaml 文件路径",
        "load_yaml_cases(path) — 读取 YAML 文件返回用例字典列表",
        "parse_test_case(raw, filename) — 校验 + 构建 TestCaseObj",
        "discover_and_parse() — 一站式发现+加载+解析",
    ]),
    ("commons/logger.py", "日志模块", [
        "## 作用",
        "统一管理框架运行时的日志输出",
        "## 使用方式",
        "from commons.logger import logger",
        "logger.info('信息') / logger.error('错误')",
        "## 日志级别",
        "debug — 调试信息（请求详情、参数）",
        "info — 关键流程（用例开始/结束）",
        "warning — 非致命警告",
        "error — 错误信息（请求失败等）",
        "## 日志文件",
        "位置: logs/YYYY-MM-DD.log，自动轮转 10MB",
    ]),
]

for idx, (mod_name, subtitle, bullets) in enumerate(module_slides):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_section_header(slide, f"模块详解 {idx+1} {mod_name}", subtitle)
    add_shape_bg(slide, CLR_PRIMARY)
    add_bullet_slide(slide, f"模块详解 {idx+1} {mod_name}", bullets)

# ====== Slide 21: debug_talk.py ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "热加载函数库 debug_talk.py", "DYNAMIC FUNCTIONS")
add_shape_bg(slide, CLR_PRIMARY)
add_two_column(slide, "debug_talk.py",
    "随机数据生成",
    ["random_number(min, max) — 随机整数",
     "random_phone() — 随机手机号",
     "random_string(length) — 随机字符串",
     "uuid_str() / uuid_with_hyphen() — UUID"],
    "日期时间",
    ["current_date(fmt) — 当前日期",
     "timestamp() / timestamp_ms() — 时间戳",
     "format_time(days, fmt) — 偏移日期",
     "get_extract_data(key) — 读取 extract.yaml"])
note2 = "YAML 中调用: username: ${random_phone()}  |  date: ${format_time(7)}  |  id: ${uuid_str()}"
add_info_box(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.5), note2)
add_page_number(slide, 21, TOTAL)

# ====== Slide 22: run.py ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "运行入口 run.py + pytest.ini", "ENTRY POINT")
add_shape_bg(slide, CLR_PRIMARY)
code_run = "# 使用方式\npython run.py                      # 默认 test 环境\npython run.py --env prod           # 生产环境\npython run.py --env dev -m smoke   # 开发环境 + 冒烟测试\n\n# 支持的参数\n--env {dev,test,prod}   环境选择\n--alluredir ./temps     Allure 结果目录\n-m, --mark             标记表达式（如 smoke）\n\n# 内部行为\n1. 解析参数 -> 2. 切换环境 (set_env)\n3. 执行 pytest -> 4. 生成 Allure 报告"
add_code_block(slide, Inches(0.5), Inches(1.6), Inches(6), Inches(5), code_run, font_size=11)
info_pytest = """## pytest.ini 关键配置
addopts = -vs --alluredir=./temps
testpaths = ./testcases
python_files = test_all_case.py
markers = smoke:冒烟测试

## conftest.py fixtures
setup_global_environment (session)
  - 清空 extract.yaml
  - 记录环境信息

exe_sql_fixture (function)
  - 预留的数据库前后置"""
add_info_box(slide, Inches(6.8), Inches(1.6), Inches(6), Inches(5), info_pytest)
add_page_number(slide, 22, TOTAL)

# ====== Slide 23: Pytest 集成 ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "Pytest 集成 test_all_case.py", "PYTEST INTEGRATION")
add_shape_bg(slide, CLR_PRIMARY)
bullets = [
    "## 作用 — 连接 YAML 用例和 pytest",
    "扫描 testcases/**/*.yaml，为每个文件动态创建 pytest 测试方法",
    "## 核心逻辑",
    "1. discover: glob('**/*.yaml') 获取所有 YAML 文件路径",
    "2. create: 对每个文件调用 creat_testcase() 生成测试函数",
    "3. inject: setattr(TestAllCase, test_name, func) 注入到测试类",
    "4. parametrize: 如果 YAML 有 parametrize，自动展开",
    "5. allure: 执行时添加 allure.dynamic.feature/story/title",
    "## 处理单个用例",
    "_run_single_case(case, path) — 校验 -> 执行",
    "校验失败 -> pytest.skip() 跳过用例",
    "执行失败 -> AssertionError -> pytest 标记失败",
]
add_bullet_slide(slide, "Pytest 集成", bullets)

# ====== Slide 24: 完整执行流程 ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "完整执行流程", "END-TO-END FLOW")
add_shape_bg(slide, CLR_PRIMARY)
flow = [
    ("Step 1", "用户执行\npython run.py", CLR_SECONDARY),
    ("Step 2", "初始化\n加载配置+清空extract.yaml", CLR_PRIMARY),
    ("Step 3", "发现用例\n扫描YAML，动态创建测试函数", RGBColor(0x00, 0x96, 0x88)),
    ("Step 4", "解析校验\nmodel_util 校验YAML格式", CLR_ORANGE),
    ("Step 5", "发送请求\nrequest_util 替换变量后发HTTP", CLR_ACCENT),
    ("Step 6", "提取数据\nextract_util 提取并写入yaml", RGBColor(0x7B, 0x1F, 0xA2)),
    ("Step 7", "断言验证\nassert_util 执行equals/contains", RGBColor(0x00, 0x96, 0x88)),
    ("Step 8", "生成报告\nAllure HTML 报告", CLR_PRIMARY),
]
for i, (step, desc, color) in enumerate(flow):
    row = i // 4
    col = i % 4
    left = Inches(0.5 + col * 3.2)
    top = Inches(1.6 + row * 2.8)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(2.8), Inches(2.3))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)
    card.line.color.rgb = color
    card.line.width = Pt(2)
    txBox = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.15), Inches(2.5), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = step
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = color
    txBox2 = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.5), Inches(2.5), Inches(1.7))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = desc
    p2.font.size = Pt(12)
    p2.font.color.rgb = CLR_DARK
    p2.line_spacing = Pt(18)
    if col < 3:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left+Inches(2.8), top+Inches(1.0), Inches(0.4), Inches(0.3))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = color
        arrow.line.fill.background()
note = "全流程自动，用户只需关注 Step 1（写 YAML）和看结果"
add_info_box(slide, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.6), note)
add_page_number(slide, 24, TOTAL)

# ====== Slide 25: 手把手教程 ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "编写第一个测试用例", "HANDS-ON TUTORIAL")
add_shape_bg(slide, CLR_PRIMARY)
steps_data = [
    ("Step 1", "创建文件", "testcases/my_test/my_first_test.yaml"),
    ("Step 2", "写用例", "复制模板，修改 method/url/params"),
    ("Step 3", "配环境", "run.env 加 base_url / config.yaml 配 URL"),
    ("Step 4", "运行", "python run.py --env test"),
    ("Step 5", "加断言", "加 validate / extract 丰富场景"),
]
for i, (step, title, desc) in enumerate(steps_data):
    row = i // 3
    col = i % 3
    left = Inches(0.5 + col * 4.2)
    top = Inches(1.6 + row * 1.5)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(3.8), Inches(1.3))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(0xF0, 0xF4, 0xFF)
    card.line.color.rgb = CLR_PRIMARY
    card.line.width = Pt(1.5)
    txBox = slide.shapes.add_textbox(left+Inches(0.2), top+Inches(0.1), Inches(3.4), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{step}: {title}"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = CLR_PRIMARY
    txBox2 = slide.shapes.add_textbox(left+Inches(0.2), top+Inches(0.4), Inches(3.4), Inches(0.8))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = desc
    p2.font.size = Pt(11)
    p2.font.color.rgb = CLR_DARK

code_tpl = """- feature: 我的模块
  story: 我的接口
  title: 验证接口返回成功
  request:
    method: get
    url: ${env(base_url)}/api/example
    params:
      key: value
  validate:
    equals:
      预期状态码200: [200, status_code]"""
add_code_block(slide, Inches(0.5), Inches(4.8), Inches(7), Inches(2.3), code_tpl, font_size=10)
tip = "写完 YAML 直接运行 python run.py，新文件自动被发现和执行"
add_info_box(slide, Inches(8), Inches(4.8), Inches(5), Inches(2.3), tip)
add_page_number(slide, 25, TOTAL)

# ====== Slide 26: 场景案例 ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "场景案例：接口关联", "REAL-WORLD EXAMPLE")
add_shape_bg(slide, CLR_PRIMARY)
code_a = """# 用例1: 获取 token
- feature: 鉴权模块
  story: 获取 access_token
  title: 获取 token 成功
  request:
    method: get
    url: ${env(base_url)}/cgi-bin/token
    params:
      appid: aaa
      secret: bbb
  extract:
    access_token: [json, "$.access_token", 0]
  validate:
    equals:
      状态码200: [200, status_code]"""
code_b = """# 用例2: 使用 token
- feature: 用户模块
  story: 获取用户信息
  title: 携带 token 获取用户信息
  request:
    method: get
    url: ${env(base_url)}/api/user/info
    params:
      token: ${read_yaml(access_token)}
  validate:
    equals:
      状态码200: [200, status_code]"""
add_code_block(slide, Inches(0.5), Inches(1.6), Inches(6), Inches(3.3), code_a, font_size=9)
add_code_block(slide, Inches(6.8), Inches(1.6), Inches(6), Inches(3.3), code_b, font_size=9)
note = "用例1 提取 access_token -> 自动写入 extract.yaml\n用例2 通过 ${read_yaml(access_token)} 读取 -> 完成接口关联\n两个 yaml 文件可以独立编写，框架自动维护变量传递"
add_info_box(slide, Inches(0.5), Inches(5.3), Inches(12.3), Inches(1.8), note)
add_page_number(slide, 26, TOTAL)

# ====== Slide 27: Allure 报告 ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "Allure 测试报告", "TEST REPORT")
add_shape_bg(slide, CLR_PRIMARY)
bullets = [
    "## 什么是 Allure？",
    "开源测试报告框架，将测试结果渲染为美观的 HTML 网页",
    "## 报告中能看到的",
    "测试概览 — 总用例数、通过数、失败数、耗时",
    "分类查看 — 按 feature/story 分层（对应 YAML 中的字段）",
    "失败详情 — 断言失败时显示期望值 vs 实际值",
    "## 生成方式",
    "python run.py → 自动生成 Allure 报告到 reports/ 目录",
    "也可用 allure serve ./temps 启动本地服务预览",
    "## YAML 与 Allure 对应",
    "YAML feature → Allure epic（一级分组：功能模块）",
    "YAML story → Allure feature（二级分组：接口名称）",
    "YAML title → Allure story（三级标题：用例标题）",
]
add_bullet_slide(slide, "Allure 测试报告", bullets)

# ====== Slide 28: FAQ ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "常见问题 FAQ", "FREQUENTLY ASKED QUESTIONS")
add_shape_bg(slide, CLR_PRIMARY)
qas = [
    ("Q: 我不懂 Python，能用吗？", "可以！写 YAML 即可。框架规范.txt 详细说明了 YAML 怎么写，复制模板改接口信息即可。"),
    ("Q: 登录 token 怎么办？", "先写获取 token 的用例（配置 extract），token 自动保存到 extract.yaml。后续用 ${read_yaml(token)} 引用。"),
    ("Q: 多环境怎么切换？", "config.yaml 配好各环境 base_url，运行 python run.py --env prod 即可。"),
    ("Q: CSV 数据驱动怎么用？", "YAML 的 parametrize 写 CSV 路径，第一行是参数名，后续每行是一组测试数据。"),
    ("Q: 测试失败怎么排查？", "1. 看控制台错误 2. 看 logs/ 日志文件 3. 看 Allure 报告的失败详情"),
]
for i, (q, a) in enumerate(qas):
    top = Inches(1.6 + i * 1.1)
    txBox = slide.shapes.add_textbox(Inches(0.8), top, Inches(11.5), Inches(0.35))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = q
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = CLR_PRIMARY
    txBox2 = slide.shapes.add_textbox(Inches(1.0), top + Inches(0.35), Inches(11.3), Inches(0.6))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = a
    p2.font.size = Pt(11)
    p2.font.color.rgb = CLR_GRAY
add_page_number(slide, 28, TOTAL)

# ====== Slide 29: 模块关系图 ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_section_header(slide, "模块关系图", "MODULE RELATIONSHIP")
add_shape_bg(slide, CLR_PRIMARY)
table_data = [
    ("模块", "依赖哪些模块", "被哪些模块依赖"),
    ("settings.py", "logger, yaml, os, dotenv", "所有其他模块"),
    ("yaml_util.py", "settings (PROJECT_ROOT)", "model_util, case_util, extract_util, debug_talk"),
    ("model_util.py", "yaml_util, Cerberus", "main_util, case_util, test_all_case"),
    ("ddt_util.py", "settings (PROJECT_ROOT)", "test_all_case"),
    ("request_util.py", "settings, logger, yaml_util", "main_util"),
    ("extract_util.py", "logger, yaml_util, jsonpath", "main_util"),
    ("assert_util.py", "jsonpath", "main_util"),
    ("main_util.py", "request_util, extract_util, assert_util", "test_all_case"),
    ("case_util.py", "yaml_util, model_util", "（一站式接口）"),
    ("test_all_case.py", "main_util, model_util, ddt_util", "pytest 收集执行"),
]
rows, cols = len(table_data), 3
table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.6), Inches(12.3), Inches(5.5)).table
for i, (c0, c1, c2) in enumerate(table_data):
    bold = (i == 0)
    c = CLR_WHITE if i == 0 else CLR_DARK
    set_cell_text(table.cell(i, 0), c0, bold=bold, size=10, color=c)
    set_cell_text(table.cell(i, 1), c1, bold=bold, size=10, color=c)
    set_cell_text(table.cell(i, 2), c2, bold=bold, size=10, color=c)
    if i == 0:
        table.cell(i, 0).fill.solid()
        table.cell(i, 0).fill.fore_color.rgb = CLR_PRIMARY
        table.cell(i, 1).fill.solid()
        table.cell(i, 1).fill.fore_color.rgb = CLR_PRIMARY
        table.cell(i, 2).fill.solid()
        table.cell(i, 2).fill.fore_color.rgb = CLR_PRIMARY
    elif i % 2 == 0:
        table.cell(i, 0).fill.solid()
        table.cell(i, 0).fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
        table.cell(i, 1).fill.solid()
        table.cell(i, 1).fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
        table.cell(i, 2).fill.solid()
        table.cell(i, 2).fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
add_page_number(slide, 29, TOTAL)

# ====== Slide 30: 结尾 ======
slide = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = CLR_PRIMARY
bg.line.fill.background()
txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(1))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "开始使用吧！"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = CLR_WHITE
p.alignment = PP_ALIGN.CENTER
txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(11), Inches(1.5))
tf2 = txBox2.text_frame
p2 = tf2.paragraphs[0]
p2.text = "从复制一个 .yaml 文件开始，运行 python run.py\n你的第一个接口自动化用例就跑起来了！"
p2.font.size = Pt(20)
p2.font.color.rgb = RGBColor(0xBB, 0xDE, 0xFB)
p2.alignment = PP_ALIGN.CENTER
p2.line_spacing = Pt(36)
txBox3 = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11), Inches(0.5))
tf3 = txBox3.text_frame
p3 = tf3.paragraphs[0]
p3.text = "动手试一试吧！有问题随时看 logs/ 下的日志"
p3.font.size = Pt(14)
p3.font.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
p3.alignment = PP_ALIGN.CENTER
add_page_number(slide, 30, TOTAL)

# ====== 保存 ======
output_path = "/Users/huanglele/PycharmProjects/InterfaceDemo/InterfaceDemo_框架介绍.pptx"
prs.save(output_path)
print(f"PPT 已生成: {output_path}")
print(f"共 {TOTAL} 页幻灯片")
